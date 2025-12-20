from pdf2docx import Converter
from docx import Document
import os

def pdf_to_docx(input_path, output_path):
    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()

def txt_to_docx(input_path, output_path):
    doc = Document()
    with open(input_path, 'r', encoding='utf-8') as f:
        content = f.read()
        # Split content into paragraphs
        paragraphs = content.split('\n\n')
        for paragraph in paragraphs:
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
    doc.save(output_path)

def docx_to_txt(input_path, output_path):
    """Extract text from DOCX file"""
    try:
        print(f"🔍 Starting DOCX to TXT conversion: {input_path} -> {output_path}")
        
        # Check if input file exists
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
        
        doc = Document(input_path)
        text_content = []
        
        print(f"📄 Found {len(doc.paragraphs)} paragraphs in document")
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n\n'.join(text_content))
        
        print(f"✅ Successfully converted DOCX to TXT: {len(text_content)} paragraphs")
            
    except Exception as e:
        print(f"❌ DOCX to TXT conversion error: {str(e)}")
        raise Exception(f"DOCX to TXT conversion failed: {str(e)}")

def docx_to_pptx(input_path, output_path):
    """Convert DOCX content to PPTX presentation"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        # Read DOCX content
        doc = Document(input_path)
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Use first paragraph as title, or default
        first_para = next((p.text for p in doc.paragraphs if p.text.strip()), "Document Presentation")
        title.text = first_para[:50] + "..." if len(first_para) > 50 else first_para
        subtitle.text = "Converted from DOCX"
        
        # Content slides
        content_layout = prs.slide_layouts[1]  # Title and Content
        current_slide = None
        slide_count = 0
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
                
            # Create new slide every 5 paragraphs or for headings
            if slide_count == 0 or slide_count % 5 == 0:
                current_slide = prs.slides.add_slide(content_layout)
                title_shape = current_slide.shapes.title
                title_shape.text = f"Content {(slide_count // 5) + 1}"
                content_shape = current_slide.placeholders[1]
                content_shape.text = text
                slide_count += 1
            else:
                # Add to existing slide
                if current_slide and len(current_slide.placeholders) > 1:
                    content_shape = current_slide.placeholders[1]
                    content_shape.text += f"\n\n{text}"
                slide_count += 1
        
        prs.save(output_path)
        
    except Exception as e:
        raise Exception(f"DOCX to PPTX conversion failed: {str(e)}")

def txt_to_pptx(input_path, output_path):
    """Create PowerPoint from text file"""
    try:
        from pptx import Presentation
        
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Split content by slides (assuming === SLIDE X === format or double newlines)
        if "=== SLIDE" in content:
            slides_content = content.split("=== SLIDE")[1:]  # Skip first empty part
        else:
            # Split by double newlines
            slides_content = content.split('\n\n')
        
        prs = Presentation()
        
        for i, slide_content in enumerate(slides_content):
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Clean up slide content
            slide_text = slide_content.strip()
            if slide_text.startswith("==="):
                # Remove slide marker
                lines = slide_text.split('\n')[1:]
                slide_text = '\n'.join(lines).strip()
            
            if slide_text:
                lines = slide_text.split('\n')
                title = lines[0] if lines else f"Slide {i+1}"
                content = '\n'.join(lines[1:]) if len(lines) > 1 else ""
                
                # Set title
                title_shape = slide.shapes.title
                title_shape.text = title[:50] + "..." if len(title) > 50 else title
                
                # Set content
                if content and len(slide.shapes) > 1:
                    content_shape = slide.shapes.placeholders[1]
                    content_shape.text = content
        
        prs.save(output_path)
        
    except Exception as e:
        raise Exception(f"Text to PPTX conversion failed: {str(e)}")