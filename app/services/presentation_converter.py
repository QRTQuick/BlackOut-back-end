from pptx import Presentation
from pptx.util import Inches
import json
import os

def convert_presentation(input_path, output_path, target_format):
    """Convert presentation files"""
    try:
        input_ext = input_path.lower().split('.')[-1]
        target_format = target_format.lower()
        
        if input_ext == 'pptx' and target_format == 'txt':
            pptx_to_text(input_path, output_path)
        elif input_ext == 'pptx' and target_format == 'json':
            pptx_to_json(input_path, output_path)
        elif input_ext == 'txt' and target_format == 'pptx':
            text_to_pptx(input_path, output_path)
        else:
            raise ValueError(f"Conversion from {input_ext} to {target_format} not supported")
            
    except Exception as e:
        raise Exception(f"Presentation conversion failed: {str(e)}")

def pptx_to_text(input_path, output_path):
    """Extract text from PowerPoint presentation"""
    try:
        prs = Presentation(input_path)
        text_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = f"=== SLIDE {i} ===\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            
            text_content.append(slide_text)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("\n".join(text_content))
            
    except Exception as e:
        raise Exception(f"PPTX to text conversion failed: {str(e)}")

def pptx_to_json(input_path, output_path):
    """Convert PowerPoint to JSON structure"""
    try:
        prs = Presentation(input_path)
        presentation_data = {
            "title": "Extracted Presentation",
            "slides": []
        }
        
        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": i,
                "content": []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data["content"].append({
                        "type": "text",
                        "content": shape.text
                    })
            
            presentation_data["slides"].append(slide_data)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(presentation_data, f, indent=2, ensure_ascii=False)
            
    except Exception as e:
        raise Exception(f"PPTX to JSON conversion failed: {str(e)}")

def text_to_pptx(input_path, output_path):
    """Create PowerPoint from text file"""
    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Split content by slides (assuming === SLIDE X === format or double newlines)
        if "=== SLIDE" in content:
            slides_content = content.split("=== SLIDE")[1:]  # Skip first empty part
        else:
            # Split by double newlines
            slides_content = content.split('\n\n')
        
        prs = Presentation()
        
        for slide_content in slides_content:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Clean up slide content
            slide_text = slide_content.strip()
            if slide_text.startswith("==="):
                # Remove slide marker
                lines = slide_text.split('\n')[1:]
                slide_text = '\n'.join(lines).strip()
            
            if slide_text:
                lines = slide_text.split('\n')
                title = lines[0] if lines else "Slide"
                content = '\n'.join(lines[1:]) if len(lines) > 1 else ""
                
                # Set title
                title_shape = slide.shapes.title
                title_shape.text = title[:50] + "..." if len(title) > 50 else title
                
                # Set content
                if content and len(slide.shapes) > 1:
                    content_shape = slide.shapes.placeholders[1]
                    content_shape.text = content
        
        prs.save(output_path)
        
    except Exception as e:
        raise Exception(f"Text to PPTX conversion failed: {str(e)}")